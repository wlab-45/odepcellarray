'''  ppt生成，不支援動畫
from pptx import Presentation
from pptx.util import Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建一个演示文稿对象
prs = Presentation()

# 添加一个幻灯片
slide_layout = prs.slide_layouts[5]  # 使用空白幻灯片布局
slide = prs.slides.add_slide(slide_layout)

# 设置幻灯片背景颜色为黑色
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)  # 黑色

# 添加一个形状（例如矩形）
left = Cm(0)  
top = Cm(0)
width = Cm(2)  
height = Cm(19)  
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, width, height
)

# 设置形状的填充颜色
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # 白色

# 保存演示文稿
prs.save('light_bar_animation.pptx')

# 自动投影
import os
os.startfile('light_bar_animation.pptx')

'''
import os
import cv2
import re
import numpy as np
import tkinter as tk
from tkinter import simpledialog
import platform
from pptx import Presentation
from pptx.util import Inches
# 创建一个Tkinter根窗口（隐藏）
root = tk.Tk()
root.withdraw()

# 获取用户输入
rect_width = simpledialog.askinteger("输入", "请输入矩形的宽度（像素）：")
rect_height = simpledialog.askinteger("输入", "请输入矩形的高度（像素）：")
speed = simpledialog.askinteger("输入", "请输入移动速度（像素/帧）：")
end_x = simpledialog.askinteger("输入", "请输入矩形的终点位置(x)：")

# 设置图像尺寸
height, width = 480, 640

# 计算矩形的初始位置和终点位置
start_x = 0
start_y = (height - rect_height) // 2

# 打印用户输入的值
print(f"矩形宽度：{rect_width} 像素")
print(f"矩形高度：{rect_height} 像素")
print(f"移动速度：{speed} 像素/帧")

# 创建一个黑色背景的图像
image = np.zeros((height, width, 3), dtype=np.uint8)

# 设置矩形的颜色
color = (255, 255, 255)  # BGR格式的白色

# 设置帧率（单位：帧/秒）
fps = 30
frame_time = int(1000 / fps)  # 每帧的时间（毫秒）

# 创建视频写入对象
video_file = f"animation_{rect_width}&{speed}(pixel_per_frame).mp4"
fourcc = cv2.VideoWriter_fourcc(*'mp4v')
out = cv2.VideoWriter(video_file, fourcc, fps, (width, height))
video_file = f"animation_{rect_width}&{speed}(pixel_per_frame).mp4"
safe_video_file = re.sub(r'[&<>:"/\\|?*]', '_', video_file)
print(f"新文件名: {safe_video_file}")
frame_count = 0
# 动画循环
x = start_x
while x <= end_x:
    # 重新绘制黑色背景
    image[:] = (0, 0, 0)
    
    # 绘制矩形
    top_left = (x, start_y)
    bottom_right = (x + rect_width, start_y + rect_height)
    cv2.rectangle(image, top_left, bottom_right, color, -1)
    
    # 写入视频帧
    out.write(image)
    
    # 更新矩形位置
    x += speed
    frame_count += 1

# 释放视频写入对象
out = cv2.VideoWriter(safe_video_file, fourcc, fps, (width, height))

print(f"动画秒数 = {frame_count * frame_time / 1000} 秒")

# # 检查视频文件是否存在
# if not os.path.exists(video_file):
#     print(f"文件 {video_file} 不存在")
# else:
#     # 根据操作系统选择合适的命令播放视频文件
#     if platform.system() == 'Windows':
#         os.startfile(video_file)
#     elif platform.system() == 'Darwin':  # macOS
#         os.system(f'open {video_file}')
#     else:  # Linux
#         os.system(f'xdg-open {video_file}')


# 创建一个 PPT 文件
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])  # 添加空白幻灯片

# 插入影片到幻灯片
video_path = video_file  # 影片路径
left = top = Inches(1)  # 影片位置
height = Inches(4)  # 影片高度（根据需要调整）

# 设置影片宽度，假设宽高比是 16:9
aspect_ratio = 16 / 9 
width = Inches(6)  # 根据高度计算宽度
# 添加影片框
movie = slide.shapes.add_movie(safe_video_file, left, top, width=width, height=height, poster_frame_image=None)

# 保存 PPT 文件
ppt_file = f"animation_{rect_width}&{speed}.pptx"
prs.save(ppt_file)
print(f"PPT 文件已保存为 {ppt_file}")